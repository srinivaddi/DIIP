import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def apply_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_slide_header(slide, title_text, color):
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(8.5), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Arial"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = color

def main():
    prs = Presentation()
    # Change slide width and height to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Custom color palette (DIIP Advisor Sleek Theme)
    bg_color = RGBColor(10, 15, 30)        # Dark Indigo Navy
    primary_color = RGBColor(56, 189, 248) # Light Blue / Sky (representing trust and advice)
    body_color = RGBColor(226, 232, 240)   # Slate-200
    accent_color = RGBColor(16, 185, 129)  # Emerald Accent
    
    blank_layout = prs.slide_layouts[6]
    
    # ----------------------------------------------------
    # SLIDE 1: Title Slide
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, bg_color)
    
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "DIIP FOR WEALTH ADVISORS"
    p.font.name = "Arial"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = primary_color
    p.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "Driving Client Engagement & Smarter Asset Allocation"
    p2.font.name = "Arial"
    p2.font.size = Pt(28)
    p2.font.color.rgb = accent_color
    p2.alignment = PP_ALIGN.LEFT
    
    p3 = tf.add_paragraph()
    p3.text = "\nIntegrating Institutional Research and Crowding Safeguards into Advisory Workflows"
    p3.font.name = "Arial"
    p3.font.size = Pt(16)
    p3.font.color.rgb = body_color
    p3.alignment = PP_ALIGN.LEFT

    # ----------------------------------------------------
    # SLIDE 2: Key Advisor Challenges
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    apply_background(slide2, bg_color)
    add_slide_header(slide2, "Key Challenges in Wealth Advisory", primary_color)
    
    box2 = slide2.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(5))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    
    bullets2 = [
        ("Managing Client Anxiety", "Clients frequently panic during short-term market rotations. Advisors spend valuable hours explaining complex macro shifts without interactive visual summaries."),
        ("Translating Professional Outlooks", "Explaining dense multi-page strategy reports from major investment banks (Morgan Stanley, Goldman Sachs) to a retail client is difficult and time-consuming."),
        ("Protecting Clients from Crowded Trades", "Advisors need a rigorous, data-driven framework to justify avoiding highly popular, over-leveraged stock peaks (retail FOMO risks).")
    ]
    
    for title, desc in bullets2:
        p_t = tf2.add_paragraph() if tf2.text else tf2.paragraphs[0]
        p_t.text = f"•  {title}"
        p_t.font.name = "Arial"
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = accent_color
        
        p_d = tf2.add_paragraph()
        p_d.text = f"    {desc}\n"
        p_d.font.name = "Arial"
        p_d.font.size = Pt(15)
        p_d.font.color.rgb = body_color

    # ----------------------------------------------------
    # SLIDE 3: The Advisor Console
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    apply_background(slide3, bg_color)
    add_slide_header(slide3, "The DIIP Advisor Console", primary_color)
    
    box3 = slide3.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(5))
    tf3 = box3.text_frame
    tf3.word_wrap = True
    
    bullets3 = [
        ("Consolidated Research Hub", "Instantly accesses real-time thematic commentary and strategic releases across 6 institutions in a single dashboard, saving hours of manual parsing."),
        ("Dynamic Multi-Desk Filters", "Allows advisors to filter alerts and research dynamically by target desk views: Macro Strategy, Equity Selection, or Portfolio Baskets."),
        ("Thematic Conviction Maps", "Provides clear institutional consensus status and safety scores to validate recommendations during client portfolio reviews.")
    ]
    
    for title, desc in bullets3:
        p_t = tf3.add_paragraph() if tf3.text else tf3.paragraphs[0]
        p_t.text = f"•  {title}"
        p_t.font.name = "Arial"
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = accent_color
        
        p_d = tf3.add_paragraph()
        p_d.text = f"    {desc}\n"
        p_d.font.name = "Arial"
        p_d.font.size = Pt(15)
        p_d.font.color.rgb = body_color

    # ----------------------------------------------------
    # SLIDE 4: Interactive Client Communication (Layman Mode)
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    apply_background(slide4, bg_color)
    add_slide_header(slide4, "Interactive Client Communication", primary_color)
    
    box4 = slide4.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(5))
    tf4 = box4.text_frame
    tf4.word_wrap = True
    
    bullets4 = [
        ("One-Click 'Layman Mode' Toggle", "Enables the advisor to show the screen directly to clients, dynamically translating technical metrics into plain-English guide terms (e.g. 'Options Call Skew' -> 'Speculative Bets')."),
        ("Dynamic Forecast Horizons", "Allows toggle guides for 30-Day, 6-Month, and 1-Year targets. Helps clients visualize short-term noise vs. structural long-term retirement value."),
        ("Actionable Threat Alerts", "Presents visual threat-level bars and flashing indicators. Explains to clients exactly how real-time news (like China export caps) affects their specific allocations.")
    ]
    
    for title, desc in bullets4:
        p_t = tf4.add_paragraph() if tf4.text else tf4.paragraphs[0]
        p_t.text = f"•  {title}"
        p_t.font.name = "Arial"
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = accent_color
        
        p_d = tf4.add_paragraph()
        p_d.text = f"    {desc}\n"
        p_d.font.name = "Arial"
        p_d.font.size = Pt(15)
        p_d.font.color.rgb = body_color

    # ----------------------------------------------------
    # SLIDE 5: Safe Portfolio Allocation
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    apply_background(slide5, bg_color)
    add_slide_header(slide5, "Risk Management & Safe Allocations", primary_color)
    
    box5 = slide5.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(5))
    tf5 = box5.text_frame
    tf5.word_wrap = True
    
    bullets5 = [
        ("Crowding Reality Safeguard", "Queries live options skew and positioning data to identify over-popular trades. Warns advisors and clients with clear 'Avoid/Red Flag' indications before entering high-risk allocations."),
        ("Consolidated Batch Price Ingestion", "Calculates live momentum dynamically using Yahoo Finance batch queries, backed by resilient Stale-While-Revalidate (SWR) recovery channels."),
        ("Thematic Buy/Sell/Hold Guide", "Maintains an active list of 18 tickers (spanning Equities, ETFs, and Mutual Funds) divided into clear green, yellow, and red action columns.")
    ]
    
    for title, desc in bullets5:
        p_t = tf5.add_paragraph() if tf5.text else tf5.paragraphs[0]
        p_t.text = f"•  {title}"
        p_t.font.name = "Arial"
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = accent_color
        
        p_d = tf5.add_paragraph()
        p_d.text = f"    {desc}\n"
        p_d.font.name = "Arial"
        p_d.font.size = Pt(15)
        p_d.font.color.rgb = body_color
        
    prs.save("diip_advisor_presentation.pptx")
    print("PowerPoint presentation successfully saved as 'diip_advisor_presentation.pptx'")

if __name__ == "__main__":
    main()
