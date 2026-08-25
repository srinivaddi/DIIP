import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def apply_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_slide_header(slide, title_text, color):
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(8.5), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Arial"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = color

def main():
    prs = Presentation()
    # Change slide width and height to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Custom color palette (DIIP Sleek Theme)
    bg_color = RGBColor(11, 13, 25)       # Dark Slate #0b0d19
    primary_color = RGBColor(16, 185, 129) # Emerald #10b981
    body_color = RGBColor(226, 232, 240)   # Slate-200 #e2e8f0
    accent_color = RGBColor(45, 212, 191)  # Teal #2dd4bf
    
    blank_layout = prs.slide_layouts[6]
    
    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Sleek Dark Theme)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, bg_color)
    
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "DIIP ENGINE"
    p.font.name = "Arial"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = primary_color
    p.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "Digital Institutional Intelligence Platform"
    p2.font.name = "Arial"
    p2.font.size = Pt(28)
    p2.font.color.rgb = accent_color
    p2.alignment = PP_ALIGN.LEFT
    
    p3 = tf.add_paragraph()
    p3.text = "\nDemocratizing Institutional Allocation Signals for Individual Investors"
    p3.font.name = "Arial"
    p3.font.size = Pt(16)
    p3.font.color.rgb = body_color
    p3.alignment = PP_ALIGN.LEFT

    # ----------------------------------------------------
    # SLIDE 2: The Core Problem
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    apply_background(slide2, bg_color)
    add_slide_header(slide2, "The Market Problem", primary_color)
    
    box2 = slide2.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(5))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    
    bullets2 = [
        ("Wall Street Complexity", "Retail investors are locked out of high-grade institutional research reports (Morgan Stanley, Goldman Sachs, Vanguard) due to dense, inaccessible financial jargon."),
        ("The Retail FOMO Cycle", "Without active positioning guides, individual investors routinely buy assets at the absolute peak of popularity, suffering immediate corrections due to institutional crowding."),
        ("No Actionable Focus", "Existing retail platforms provide raw news feeds without clear action triggers (Buy, Sell, or Hold), leading to analysis paralysis for the layman user.")
    ]
    
    for title, desc in bullets2:
        p_t = tf2.add_paragraph() if tf2.text else tf2.paragraphs[0]
        p_t.text = f"•  {title}"
        p_t.font.name = "Arial"
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = accent_color
        
        p_d = tf2.add_paragraph()
        p_d.text = f"    {desc}\n"
        p_d.font.name = "Arial"
        p_d.font.size = Pt(15)
        p_d.font.color.rgb = body_color

    # ----------------------------------------------------
    # SLIDE 3: The DIIP Solution
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    apply_background(slide3, bg_color)
    add_slide_header(slide3, "The DIIP Engine Solution", primary_color)
    
    box3 = slide3.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(5))
    tf3 = box3.text_frame
    tf3.word_wrap = True
    
    bullets3 = [
        ("Unified Asset Stance Signals", "Groups public equities, ETFs, and Mutual Funds into simple green, yellow, and red categories: Buy (Green Light), Hold (Neutral), Avoid (Red Flag)."),
        ("Live Multi-Scraper Ingestion", "Parses and structures commentaries dynamically from 6 leading global desks (Morgan Stanley, Vanguard, Fidelity, BlackRock, Goldman Sachs, JPMorgan)."),
        ("Automated Jargon Translation", "Includes a client-side 'Layman Mode' switch that translates professional terminology on-the-fly (e.g. 'Options Call Skew' becomes 'Speculative Bullish Bets')."),
        ("Live Quantitative Re-ranking", "Connects to Yahoo Finance Spark API in single consolidated batch calls to dynamically calculate momentum, automatically downgraded if high institutional crowding is detected.")
    ]
    
    for title, desc in bullets3:
        p_t = tf3.add_paragraph() if tf3.text else tf3.paragraphs[0]
        p_t.text = f"•  {title}"
        p_t.font.name = "Arial"
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = accent_color
        
        p_d = tf3.add_paragraph()
        p_d.text = f"    {desc}\n"
        p_d.font.name = "Arial"
        p_d.font.size = Pt(15)
        p_d.font.color.rgb = body_color

    # ----------------------------------------------------
    # SLIDE 4: Architecture & Multi-Agent Design
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    apply_background(slide4, bg_color)
    add_slide_header(slide4, "Multi-Agent System Architecture", primary_color)
    
    box4 = slide4.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(5))
    tf4 = box4.text_frame
    tf4.word_wrap = True
    
    bullets4 = [
        ("Research Agent (Crawler)", "Asynchronously fetches commentaries, PDFs, and releases from institutional desks, structuring them into clean Markdown format."),
        ("Positioning Agent (Overweight Detector)", "Monitors Options Call Skew, Short Interest Ratio, and CFTC Net Long placements to compute quantitative Crowding Scores (0-100%)."),
        ("Narrative Agent (Regime Classifier)", "Classifies monthly macroeconomic indicators (CPI, PPI inflation, Fed interest rate spreads) to detect thematic shifts over time."),
        ("Portfolio Management Agent", "Validates asset allocations, provides weight structures, and maintains a resilience channel caching Yahoo Finance spark inputs via Stale-While-Revalidate (SWR) rules.")
    ]
    
    for title, desc in bullets4:
        p_t = tf4.add_paragraph() if tf4.text else tf4.paragraphs[0]
        p_t.text = f"•  {title}"
        p_t.font.name = "Arial"
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = accent_color
        
        p_d = tf4.add_paragraph()
        p_d.text = f"    {desc}\n"
        p_d.font.name = "Arial"
        p_d.font.size = Pt(15)
        p_d.font.color.rgb = body_color

    # ----------------------------------------------------
    # SLIDE 5: Strategic Business Impact
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    apply_background(slide5, bg_color)
    add_slide_header(slide5, "Business Impact & Value Proposition", primary_color)
    
    box5 = slide5.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(5))
    tf5 = box5.text_frame
    tf5.word_wrap = True
    
    bullets5 = [
        ("Democratic Wealth Management", "Empowers retail platforms to offer institution-grade market insights to everyday traders, building high product loyalty."),
        ("Dynamic Forecast Horizons", "Provides toggles for 30-Day, 6-Month, and 1-Year outlooks, aligning recommendations to the user's specific long-term target time frames rather than short-term price noise."),
        ("Peak Risk Mitigation", "Prevents retail trading accounts from buying into over-popular, highly leveraged trades by displaying visible 'Too Crowded' warnings."),
        ("Scalable Platform Extensions", "FastAPI backend and React frontend are fully compiled, optimized, and ready for REST/WebSocket production deployments.")
    ]
    
    for title, desc in bullets5:
        p_t = tf5.add_paragraph() if tf5.text else tf5.paragraphs[0]
        p_t.text = f"•  {title}"
        p_t.font.name = "Arial"
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = accent_color
        
        p_d = tf5.add_paragraph()
        p_d.text = f"    {desc}\n"
        p_d.font.name = "Arial"
        p_d.font.size = Pt(15)
        p_d.font.color.rgb = body_color
        
    prs.save("diip_executive_presentation.pptx")
    print("PowerPoint presentation successfully saved as 'diip_executive_presentation.pptx'")

if __name__ == "__main__":
    main()
